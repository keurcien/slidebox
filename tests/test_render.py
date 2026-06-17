"""render() — python-pptx output tests.

Inspect the rendered Presentation: shape counts, names (object ids),
background fills.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.util import Emu

from slidebox import BrandTheme, Deck, render
from slidebox.types import SLIDE_H_EMU, SLIDE_W_EMU


def _shape_names(slide) -> list[str]:
    return [s.name for s in slide.shapes]


def test_deck_renders_one_slide_per_slide() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(bg="beige", object_id="s1")
        .header("a", col=1, row=1, span=(8, 2))
        .slide(bg="black", object_id="s2")
        .header("b", col=1, row=1, span=(8, 2))
    ).build()
    prs = render(deck)
    assert len(prs.slides) == 2
    assert prs.slide_width == Emu(SLIDE_W_EMU)
    assert prs.slide_height == Emu(SLIDE_H_EMU)


def test_kpi_emits_frame_and_named_parts() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .kpi(label="GMV", value="4,2", unit="M€", delta="+24%", delta_dir="up",
             col=1, row=1, span=(10, 5), object_id="hero_kpi")
    ).build()
    names = _shape_names(render(deck).slides[0])
    for suffix in ("__frame", "__label", "__value", "__delta"):
        assert f"hero_kpi{suffix}" in names


def test_object_id_becomes_shape_name() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .header("a", col=1, row=1, span=(8, 2), object_id="cover_title")
    ).build()
    assert "cover_title" in _shape_names(render(deck).slides[0])


def test_image_placeholder_draws_a_shape() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .image(placeholder_tone="nude", col=1, row=1, span=(6, 4),
               object_id="img")
    ).build()
    assert "img" in _shape_names(render(deck).slides[0])


def test_custom_theme_sets_background() -> None:
    from slidebox.types import RGB

    custom = BrandTheme(beige=RGB(1, 2, 3))
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(bg="beige", object_id="s")
        .header("a", col=1, row=1, span=(8, 2))
    ).build()
    slide = render(deck, theme=custom).slides[0]
    fill = slide.background.fill
    assert (fill.fore_color.rgb[0], fill.fore_color.rgb[1], fill.fore_color.rgb[2]) == (1, 2, 3)


def test_absolute_panel_renders_at_exact_emu() -> None:
    from slidebox import AbsoluteBox
    from slidebox.schema import PanelCard

    deck = Deck.new(title="T", object_id="t")
    sb = deck.slide(bg="white", object_id="s")
    sb._add(PanelCard(object_id="bg",
                      bbox=AbsoluteBox(x=0, y=0, w=4131600, h=5143500)))
    sb.image(placeholder_tone="nude", col=1, row=1, span=(5, 8), object_id="img")
    built = deck.build()

    shapes = {s.name: s for s in render(built).slides[0].shapes}
    panel = shapes["bg"]
    assert (panel.left, panel.top, panel.width, panel.height) == (0, 0, 4131600, 5143500)
    # Panel is drawn first, so it sits behind the image.
    names = _shape_names(render(built).slides[0])
    assert names.index("bg") < names.index("img")


def _runs(shape):
    return [(r.text, bool(r.font.bold)) for p in shape.text_frame.paragraphs for r in p.runs]


def test_body_inline_bold_splits_runs() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .body(["plain **bold** tail"], col=1, row=1, span=(6, 2), object_id="b")
    ).build()
    shapes = {s.name: s for s in render(deck).slides[0].shapes}
    assert _runs(shapes["b"]) == [("plain ", False), ("bold", True), (" tail", False)]


def test_text_color_override() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .eyebrow("X", col=1, row=1, span=(4, 1), color="#D1AE9B", object_id="e")
    ).build()
    run = render(deck).slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    assert (run.font.color.rgb[0], run.font.color.rgb[1], run.font.color.rgb[2]) == (0xD1, 0xAE, 0x9B)


def test_panel_shapes_and_outline() -> None:
    from pptx.enum.shapes import MSO_SHAPE

    from slidebox import AbsoluteBox
    from slidebox.schema import PanelCard

    deck = Deck.new(title="T", object_id="t")
    sb = deck.slide(object_id="s")
    sb._add(PanelCard(object_id="dot", shape="ellipse", fill="#D1AE9B",
                      outline="#FFF9ED", outline_pt=3,
                      bbox=AbsoluteBox(x=0, y=0, w=204000, h=204000)))
    shapes = {s.name: s for s in render(deck.build()).slides[0].shapes}
    dot = shapes["dot"]
    assert dot.auto_shape_type == MSO_SHAPE.OVAL
    assert (dot.line.color.rgb[0], dot.line.color.rgb[1], dot.line.color.rgb[2]) == (0xFF, 0xF9, 0xED)


def test_text_align_center() -> None:
    from pptx.enum.text import PP_ALIGN
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .header("X", col=1, row=1, span=(8, 2), align="center", object_id="h")
    ).build()
    p = render(deck).slides[0].shapes[0].text_frame.paragraphs[0]
    assert p.alignment == PP_ALIGN.CENTER


def test_image_outline_on_placeholder() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .image(placeholder_tone="nude", col=1, row=1, span=(4, 4),
               outline="#FFFFFF", outline_pt=2.5, object_id="img")
    ).build()
    img = {s.name: s for s in render(deck).slides[0].shapes}["img"]
    assert (img.line.color.rgb[0], img.line.color.rgb[1], img.line.color.rgb[2]) == (255, 255, 255)


def test_table_renders_native_with_styling() -> None:
    from slidebox import TableCell
    from slidebox.schema import TableCard

    cells = [
        [TableCell(text="H", bold=True, fill="#F9F0E0", color="#5F6365"),
         TableCell(text="x", align="right")],
        [TableCell(text="1", fill="#F9F0E0"), TableCell(text="2", align="right")],
    ]
    deck = Deck.new(title="T", object_id="t")
    sb = deck.slide(object_id="s")
    sb._add(TableCard(object_id="tbl", cells=cells,
                      col_widths=[1000000, 2000000], row_heights=[400000, 400000],
                      border="#F1E4CD", border_pt=0.75,
                      bbox=__import__("slidebox").AbsoluteBox(
                          x=500000, y=500000, w=3000000, h=800000)))
    shapes = {s.name: s for s in render(deck.build()).slides[0].shapes}
    gf = shapes["tbl"]
    assert gf.has_table
    t = gf.table
    assert (len(t.rows), len(t.columns)) == (2, 2)
    assert t.columns[0].width == 1000000
    assert t.cell(0, 0).text == "H"
    assert t.cell(0, 0).fill.fore_color.rgb == RGBColor(0xF9, 0xF0, 0xE0)


def test_body_line_spacing_override_is_applied() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .body("★ ★ ★ ★ ★", col=1, row=1, span=(4, 1), line_spacing=1.0, object_id="r")
    ).build()
    shape = {s.name: s for s in render(deck).slides[0].shapes}["r"]
    assert shape.text_frame.paragraphs[0].line_spacing == 1.0


def test_body_default_line_spacing_is_loose() -> None:
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .body("Some paragraph copy.", col=1, row=1, span=(6, 2), object_id="b")
    ).build()
    shape = {s.name: s for s in render(deck).slides[0].shapes}["b"]
    assert shape.text_frame.paragraphs[0].line_spacing == 1.6


def test_image_crop_contain_centers_and_preserves_ratio() -> None:
    import io

    from PIL import Image

    # A wide 4:1 image placed in a square box. "contain" should shrink it to
    # fit the width and center it vertically inside the box.
    buf = io.BytesIO()
    Image.new("RGB", (400, 100), (10, 20, 30)).save(buf, format="PNG")
    buf.seek(0)
    tmp = io.BytesIO(buf.getvalue())

    import tempfile

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        f.write(tmp.getvalue())
        path = f.name

    box = 2_000_000
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .image(path=path, crop="contain", x=0, y=0, w=box, h=box, object_id="img")
    ).build()
    pic = {s.name: s for s in render(deck).slides[0].shapes}["img"]
    # 4:1 image in a square box -> displayed height is a quarter of the width.
    assert pic.width == box
    assert pic.height == box // 4
    # Centered vertically.
    assert pic.top == (box - pic.height) // 2


def test_image_crop_cover_fills_box_exactly() -> None:
    import io
    import tempfile

    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (400, 100), (10, 20, 30)).save(buf, format="PNG")
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        f.write(buf.getvalue())
        path = f.name

    box = 2_000_000
    deck = (
        Deck.new(title="T", object_id="t")
        .slide(object_id="s")
        .image(path=path, crop="cover", x=0, y=0, w=box, h=box, object_id="img")
    ).build()
    pic = {s.name: s for s in render(deck).slides[0].shapes}["img"]
    # cover fills the whole box exactly (the image is center-cropped to fit).
    assert (pic.left, pic.top, pic.width, pic.height) == (0, 0, box, box)


def test_table_rejects_ragged_rows() -> None:
    import pytest

    from slidebox import TableCell
    from slidebox.schema import TableCard
    with pytest.raises(Exception):
        TableCard(object_id="t", col_start=1, col_span=4, row_start=1, row_span=4,
                  cells=[[TableCell(text="a"), TableCell(text="b")], [TableCell(text="c")]])
