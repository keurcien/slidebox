"""Render a typed Deck into a python-pptx Presentation.

This is the translation boundary: a validated `Deck` (the source of
truth) becomes a `pptx.Presentation` in memory. Slidebox never writes to
disk or talks to a network here — `save()` and `to_google_slides()` in
`slidebox.drive` consume the Presentation this module produces.

Layout comes from `grid.cell_to_emu` (12x8 grid -> EMU bbox); typography
and colour come from the `BrandTheme`. The user never sets fonts or sizes
directly — the card type does, exactly as in the original Slides backend.

Font sizes are expressed in **real points**. (The old Slides translator
authored sizes on a 1920x1080 design canvas and scaled by 0.375 before
sending; python-pptx interprets points against the real slide, so we bake
that factor in here and store the actual rendered sizes.)
"""

from __future__ import annotations

import io
import re
import urllib.request
from collections.abc import Callable
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from slidebox.grid import cell_to_emu
from slidebox.schema import (
    BodyCard,
    Card,
    Deck,
    EyebrowCard,
    HeaderCard,
    ImageCard,
    KpiCard,
    LogoCard,
    PanelCard,
    Slide,
    SubtitleCard,
    TableCard,
)
from slidebox.theme import BrandTheme
from slidebox.types import RGB, SLIDE_H_EMU, SLIDE_W_EMU

# ---- size tables (real points) --------------------------------------
_HEADER_PT = {"h1": 24.75, "display": 39.375, "keyword": 61.875}
_KPI_VALUE_PT = {"sm": 29.25, "md": 45.0, "lg": 61.875, "xl": 82.5}
_LOGO_PT = {"sm": 13.875, "md": 18.0, "lg": 24.0}
_SUBTITLE_PT = 21.0
_BODY_PT = 19.125
_BODY_LINE_HEIGHT = 1.6
_EYEBROW_PT = 18.0
_KPI_LABEL_PT = 13.875
_KPI_DELTA_PT = 16.125
_KPI_UNIT_RATIO = 0.4
_CAPTION_PT = 12.0

_KPI_INSET_X = int(Pt(12))
_KPI_INSET_Y = int(Pt(10.5))
_KPI_BORDER_PT = 1.0
_CAPTION_GAP = int(Pt(6))
_CAPTION_H = int(Pt(18))


# ---- low-level primitives -------------------------------------------
def _rgb(c: RGB) -> RGBColor:
    return RGBColor(c.r, c.g, c.b)


def _bbox(card: Card, grid: str) -> tuple[int, int, int, int]:
    b = getattr(card, "bbox", None)
    if b is not None:
        return (b.x, b.y, b.w, b.h)
    return cell_to_emu(
        card.col_start, card.col_span, card.row_start, card.row_span, res=grid  # type: ignore[arg-type]
    )


# Text fitting. A python-pptx text box does not shrink text to fit, and
# Google Slides ignores the OOXML autofit flag on import — so we bake the
# fitted point size into the file ourselves. Width comes either from real
# font metrics (when font files are supplied via `render(deck, fonts=…)`)
# or, as a fallback, a conservative ~0.6em-per-glyph estimate.
_EMU_PER_PT = 12700
_INSET_X_PT = 14.0   # ~0.1" left + right default inset
_INSET_Y_PT = 8.0    # ~0.05" top + bottom, plus slack
_GLYPH_EM = 0.6      # avg glyph advance as a fraction of font size (errs wide)
_WIDTH_SAFETY = 0.97  # leave headroom so a borderline line doesn't wrap

FontSpec = str | dict[str, str]
Fonts = dict[str, FontSpec]


class _FontBook:
    """Resolves (family, bold, italic, size) -> a Pillow font for measuring."""

    def __init__(self, fonts: Fonts) -> None:
        self._files: dict[str, dict[str, str]] = {
            fam: ({"regular": spec} if isinstance(spec, str) else dict(spec))
            for fam, spec in fonts.items()
        }
        self.missing: set[str] = set()

    def path(self, family: str, bold: bool, italic: bool) -> str | None:
        variants = self._files.get(family)
        if not variants:
            self.missing.add(family)
            return None
        keys = []
        if bold and italic:
            keys.append("bold_italic")
        if bold:
            keys.append("bold")
        if italic:
            keys.append("italic")
        keys.append("regular")
        for k in keys:
            if k in variants:
                return variants[k]
        return next(iter(variants.values()), None)

    def pil_font(self, family: str, bold: bool, italic: bool, size_pt: float) -> Any | None:
        from PIL import ImageFont

        path = self.path(family, bold, italic)
        if path is None:
            return None
        # Pillow sizes in px; at 72 dpi 1px == 1pt, so widths come out in pt.
        return ImageFont.truetype(path, size=max(1, round(size_pt)))


def _fit(
    texts: list[str],
    w_emu: int,
    h_emu: int,
    nominal_pt: float,
    *,
    line_height: float,
    min_pt: float,
    measure: Callable[[str, float], float],
) -> float:
    """Largest size ≤ nominal at which every paragraph wraps within the box.

    `measure(text, size)` returns the rendered width of `text` in points.
    No single word may exceed the line width (prevents mid-word breaks).
    """
    w_pt = (w_emu / _EMU_PER_PT - _INSET_X_PT) * _WIDTH_SAFETY
    h_pt = h_emu / _EMU_PER_PT - _INSET_Y_PT
    if w_pt <= 0 or h_pt <= 0:
        return nominal_pt
    words_per_para = [t.split() for t in texts]
    # Step in whole points: the size we *measure* is the size we *apply*, so
    # a fractional size can't render wider than what was checked.
    size = float(int(nominal_pt))
    while size > min_pt:
        space_w = measure(" ", size)
        total_lines = 0
        widest_word = 0.0
        for words in words_per_para:
            if not words:
                total_lines += 1
                continue
            line_w = 0.0
            lines = 1
            for word in words:
                ww = measure(word, size)
                widest_word = max(widest_word, ww)
                add = ww if line_w == 0 else space_w + ww
                if line_w == 0 or line_w + add <= w_pt:
                    line_w += add
                else:
                    lines += 1
                    line_w = ww
            total_lines += lines
        if total_lines * size * line_height <= h_pt and widest_word <= w_pt:
            break
        size -= 1.0
    return max(size, min_pt)


class _Fitter:
    """Picks a font size that fits a box, using real metrics when available."""

    def __init__(self, book: _FontBook | None) -> None:
        self._book = book

    def fit(
        self,
        texts: list[str],
        w_emu: int,
        h_emu: int,
        nominal_pt: float,
        *,
        family: str,
        bold: bool = False,
        italic: bool = False,
        line_height: float = 1.2,
        min_pt: float = 7.0,
    ) -> float:
        book = self._book
        if book is None:
            def measure(text: str, size: float) -> float:
                return len(text) * _GLYPH_EM * size
        else:
            def measure(text: str, size: float) -> float:
                fnt = book.pil_font(family, bold, italic, size)
                if fnt is None:  # family not supplied — fall back to estimate
                    return len(text) * _GLYPH_EM * size
                width: float = fnt.getlength(text)
                return width

        return _fit(
            texts, w_emu, h_emu, nominal_pt,
            line_height=line_height, min_pt=min_pt, measure=measure,
        )


def _style_run(
    run: Any,
    *,
    font: str,
    size: float,
    color: RGB,
    bold: bool = False,
    italic: bool = False,
    small_caps: bool = False,
) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    if small_caps:
        # python-pptx has no small-caps property; set the OOXML attribute.
        run.font._rPr.set("cap", "small")


def _textbox(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    *,
    name: str,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return box


def _text(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    text: str,
    *,
    name: str,
    font: str,
    size: float,
    color: RGB,
    fitter: _Fitter,
    bold: bool = False,
    italic: bool = False,
    small_caps: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
    exact_size: float | None = None,
) -> Any:
    box = _textbox(slide, x, y, w, h, name=name, anchor=anchor)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    if exact_size is not None:
        size = exact_size
    else:
        size = fitter.fit(
            [text], w, h, size, family=font, bold=bold, italic=italic,
            line_height=line_spacing or 1.2,
        )
    _style_run(
        run, font=font, size=size, color=color, bold=bold, italic=italic,
        small_caps=small_caps,
    )
    return box


_MSO_SHAPE = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
}

_PP_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def _rect(
    slide: Any,
    x: int,
    y: int,
    w: int,
    h: int,
    *,
    name: str,
    fill: RGB | None,
    line_color: RGB | None = None,
    line_pt: float = 1.0,
    rounded: bool = False,
    shape: str = "rectangle",
    rotation: float = 0.0,
) -> Any:
    kind = "rounded" if (rounded and shape == "rectangle") else shape
    obj = slide.shapes.add_shape(
        _MSO_SHAPE[kind], Emu(x), Emu(y), Emu(w), Emu(h),
    )
    obj.name = name
    obj.shadow.inherit = False
    if rotation:
        obj.rotation = rotation
    if fill is None:
        obj.fill.background()
    else:
        obj.fill.solid()
        obj.fill.fore_color.rgb = _rgb(fill)
    if line_color is None:
        obj.line.fill.background()
    else:
        obj.line.color.rgb = _rgb(line_color)
        obj.line.width = Pt(line_pt)
    return obj


# ---- per-card emitters ----------------------------------------------
def _emit_header(c: HeaderCard, slide: Slide, theme: BrandTheme, page: Any,
                 fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    color = RGB.from_hex(c.color) if c.color else theme.text_on(slide.background)
    _text(
        page, x, y, w, h, c.text, name=c.object_id, fitter=fitter,
        font=theme.serif_family, size=_HEADER_PT[c.size], bold=True,
        color=color, exact_size=c.size_pt, align=_PP_ALIGN[c.align],
    )


def _emit_subtitle(c: SubtitleCard, slide: Slide, theme: BrandTheme, page: Any,
                   fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    _text(
        page, x, y, w, h, c.text, name=c.object_id, fitter=fitter,
        font=theme.sans_family, size=_SUBTITLE_PT, color=theme.grey_700,
    )


def _emit_eyebrow(c: EyebrowCard, slide: Slide, theme: BrandTheme, page: Any,
                  fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    sans = c.variant == "sans"
    color = RGB.from_hex(c.color) if c.color else theme.grey_500
    _text(
        page, x, y, w, h, c.text, name=c.object_id, fitter=fitter,
        font=theme.sans_family if sans else theme.serif_family,
        size=_EYEBROW_PT, italic=not sans, color=color,
        exact_size=c.size_pt, align=_PP_ALIGN[c.align],
    )


_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


def _bold_runs(text: str) -> list[tuple[str, bool]]:
    """Split `**emphasis**` markup into (segment, bold) runs."""
    runs: list[tuple[str, bool]] = []
    i = 0
    for m in _BOLD_RE.finditer(text):
        if m.start() > i:
            runs.append((text[i:m.start()], False))
        runs.append((m.group(1), True))
        i = m.end()
    if i < len(text):
        runs.append((text[i:], False))
    return runs or [(text, False)]


def _emit_body(c: BodyCard, slide: Slide, theme: BrandTheme, page: Any,
               fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    box = _textbox(page, x, y, w, h, name=c.object_id)
    tf = box.text_frame
    if c.color:
        color = RGB.from_hex(c.color)
    elif c.tone == "muted":
        color = theme.grey_700
    else:
        color = theme.text_on(slide.background)
    line_height = c.line_spacing if c.line_spacing is not None else _BODY_LINE_HEIGHT
    # Measure with markup stripped so widths reflect the real glyphs.
    plain = [_BOLD_RE.sub(r"\1", p) for p in c.paragraphs]
    if c.size_pt is not None:
        size = c.size_pt
    else:
        size = fitter.fit(plain, w, h, _BODY_PT, family=theme.sans_family,
                          line_height=line_height)
    strong = set(c.strong)
    for i, para in enumerate(c.paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_height
        p.alignment = _PP_ALIGN[c.align]
        for seg, seg_bold in _bold_runs(para):
            run = p.add_run()
            run.text = seg
            _style_run(run, font=theme.sans_family, size=size, color=color,
                       bold=seg_bold or i in strong)


def _emit_kpi(c: KpiCard, slide: Slide, theme: BrandTheme, page: Any,
              fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    color = theme.text_on(slide.background)

    # Hairline-bordered, transparent frame across the whole cell.
    _rect(
        page, x, y, w, h, name=f"{c.object_id}__frame",
        fill=None, line_color=theme.kpi_border_for(slide.background),
        line_pt=_KPI_BORDER_PT,
    )

    x += _KPI_INSET_X
    y += _KPI_INSET_Y
    w -= 2 * _KPI_INSET_X
    h -= 2 * _KPI_INSET_Y

    label_h = int(h * 0.18) if c.label else 0
    delta_h = int(h * 0.16) if c.delta else 0
    value_h = h - label_h - delta_h
    cy = y

    if c.label:
        _text(
            page, x, cy, w, label_h, c.label, name=f"{c.object_id}__label",
            fitter=fitter, font=theme.sans_family, size=_KPI_LABEL_PT,
            color=theme.grey_500, small_caps=True,
        )
        cy += label_h

    vbox = _textbox(page, x, cy, w, value_h, name=f"{c.object_id}__value",
                    anchor=MSO_ANCHOR.MIDDLE)
    vp = vbox.text_frame.paragraphs[0]
    vrun = vp.add_run()
    vrun.text = c.value
    # Shrink the hero number (value + unit) to fit the value band.
    value_line = c.value if not c.unit else f"{c.value} {c.unit}"
    vsize = fitter.fit([value_line], w, value_h, _KPI_VALUE_PT[c.size],
                       family=theme.serif_family, bold=True)
    _style_run(vrun, font=theme.serif_family, size=vsize, bold=True, color=color)
    if c.unit:
        urun = vp.add_run()
        urun.text = f" {c.unit}"
        _style_run(urun, font=theme.sans_family,
                   size=vsize * _KPI_UNIT_RATIO, color=theme.grey_700)
    cy += value_h

    if c.delta:
        _text(
            page, x, cy, w, delta_h, c.delta, name=f"{c.object_id}__delta",
            fitter=fitter, font=theme.sans_family, size=_KPI_DELTA_PT,
            color=theme.delta_color(c.delta_dir),
        )


def _fetch_image(c: ImageCard) -> io.BytesIO | str | None:
    """Resolve an image card to something add_picture accepts, or None."""
    if c.source_url:
        if c.source_url.startswith(("http://", "https://")):
            try:
                with urllib.request.urlopen(c.source_url, timeout=10) as resp:
                    return io.BytesIO(resp.read())
            except Exception:
                return None
        return c.source_url  # local path
    return None  # drive_file_id needs auth — handled as placeholder in v1


def _crop_cover(src: io.BytesIO | str, w: int, h: int) -> io.BytesIO | str:
    """Center-crop `src` to the box's aspect ratio (CSS `object-fit: cover`).

    Returns a PNG/JPEG buffer of the cropped image; on any failure (unreadable
    image, missing codec) returns `src` unchanged so the box still fills.
    """
    try:
        from PIL import Image

        img = Image.open(src)
        iw, ih = img.size
        if iw <= 0 or ih <= 0 or h <= 0:
            return src
        target = w / h
        ratio = iw / ih
        if ratio > target:  # image too wide — trim left/right
            new_w = max(1, round(ih * target))
            x0 = (iw - new_w) // 2
            box = (x0, 0, x0 + new_w, ih)
        else:  # image too tall — trim top/bottom
            new_h = max(1, round(iw / target))
            y0 = (ih - new_h) // 2
            box = (0, y0, iw, y0 + new_h)
        cropped = img.crop(box)
        fmt = (img.format or "PNG").upper()
        if fmt in ("JPEG", "JPG") and cropped.mode not in ("RGB", "L"):
            cropped = cropped.convert("RGB")
        buf = io.BytesIO()
        cropped.save(buf, format=fmt)
        buf.seek(0)
        return buf
    except Exception:
        if hasattr(src, "seek"):
            src.seek(0)
        return src


def _contain_box(src: io.BytesIO | str, x: int, y: int, w: int, h: int
                 ) -> tuple[int, int, int, int]:
    """Largest box preserving the image's aspect ratio, centered in (x,y,w,h).

    For CSS `object-fit: contain`. Returns the original box on any failure.
    """
    try:
        from PIL import Image

        img = Image.open(src)
        iw, ih = img.size
        if hasattr(src, "seek"):
            src.seek(0)
        if iw <= 0 or ih <= 0:
            return x, y, w, h
        scale = min(w / iw, h / ih)
        dw = max(1, round(iw * scale))
        dh = max(1, round(ih * scale))
        return x + (w - dw) // 2, y + (h - dh) // 2, dw, dh
    except Exception:
        if hasattr(src, "seek"):
            src.seek(0)
        return x, y, w, h


def _emit_image(c: ImageCard, slide: Slide, theme: BrandTheme, page: Any,
                fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    outline = RGB.from_hex(c.outline) if c.outline else None
    src = _fetch_image(c)
    if src is not None:
        px, py, pw, ph = x, y, w, h
        if c.crop == "cover":
            src = _crop_cover(src, w, h)
        elif c.crop == "contain":
            px, py, pw, ph = _contain_box(src, x, y, w, h)
        try:
            pic = page.shapes.add_picture(src, Emu(px), Emu(py), Emu(pw), Emu(ph))
            pic.name = c.object_id
            if c.rotation:
                pic.rotation = c.rotation
            if outline is not None:
                pic.line.color.rgb = _rgb(outline)
                pic.line.width = Pt(c.outline_pt)
        except Exception:
            src = None
    if src is None:
        _rect(page, x, y, w, h, name=c.object_id, fill=theme.nude,
              rounded=c.rounded, line_color=outline, line_pt=c.outline_pt,
              rotation=c.rotation)

    if c.caption:
        _text(
            page, x, y + h + _CAPTION_GAP, w, _CAPTION_H, c.caption,
            name=f"{c.object_id}__caption", fitter=fitter,
            font=theme.sans_family, size=_CAPTION_PT, color=theme.grey_500,
        )


def _emit_panel(c: PanelCard, slide: Slide, theme: BrandTheme, page: Any,
                fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    fill = RGB.from_hex(c.fill) if c.fill else theme.background_rgb(c.tone)
    line = RGB.from_hex(c.outline) if c.outline else None
    _rect(page, x, y, w, h, name=c.object_id, fill=fill, rounded=c.rounded,
          shape=c.shape, rotation=c.rotation,
          line_color=line, line_pt=c.outline_pt)


def _cell_borders(tc: Any, color: RGB, pt: float) -> None:
    """Set a uniform solid border on all four edges of a table cell."""
    tcPr = tc._tc.get_or_add_tcPr()
    width = str(int(pt * _EMU_PER_PT))
    hexc = f"{color.r:02X}{color.g:02X}{color.b:02X}"
    # Schema order: lnL, lnR, lnT, lnB come first in tcPr — insert at the
    # front in reverse so they end up correctly ordered before any fill.
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
        ln = tcPr.makeelement(qn(tag),
                              {"w": width, "cap": "flat", "cmpd": "sng", "algn": "ctr"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": hexc})
        fill.append(clr)
        ln.append(fill)
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "solid"}))
        tcPr.insert(0, ln)


def _emit_table(c: TableCard, slide: Slide, theme: BrandTheme, page: Any,
                fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    nrow, ncol = len(c.cells), len(c.cells[0])
    gf = page.shapes.add_table(nrow, ncol, Emu(x), Emu(y), Emu(w), Emu(h))
    gf.name = c.object_id
    tbl = gf.table
    # Drop the inherited table style's special formatting — we set every
    # cell's fill and border explicitly.
    tbl.first_row = tbl.last_row = tbl.first_col = tbl.last_col = False
    tbl.horz_banding = tbl.vert_banding = False
    if c.col_widths:
        for i, cw in enumerate(c.col_widths):
            tbl.columns[i].width = Emu(cw)
    if c.row_heights:
        for i, rh in enumerate(c.row_heights):
            tbl.rows[i].height = Emu(rh)

    font = c.font or theme.sans_family
    border = RGB.from_hex(c.border) if c.border else None
    default_color = theme.text_on(slide.background)
    for ri, row in enumerate(c.cells):
        for ci, cell in enumerate(row):
            tc = tbl.cell(ri, ci)
            tc.vertical_anchor = MSO_ANCHOR.MIDDLE
            tc.margin_left = tc.margin_right = Pt(5)
            tc.margin_top = tc.margin_bottom = Pt(1)
            # A cell without an explicit fill takes the slide background, so
            # it blends in. (Leaving it "no fill" would expose the inherited
            # PowerPoint table style — a blue accent — on conversion.)
            cell_fill = (RGB.from_hex(cell.fill) if cell.fill
                         else theme.background_rgb(slide.background))
            tc.fill.solid()
            tc.fill.fore_color.rgb = _rgb(cell_fill)
            color = RGB.from_hex(cell.color) if cell.color else default_color
            size = cell.size_pt if cell.size_pt is not None else _BODY_PT
            tf = tc.text_frame
            tf.word_wrap = True
            for li, line in enumerate(cell.text.split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = _PP_ALIGN[cell.align]
                run = p.add_run()
                run.text = line
                _style_run(run, font=font, size=size, color=color, bold=cell.bold)
            if border is not None:
                _cell_borders(tc, border, c.border_pt)


def _emit_logo(c: LogoCard, slide: Slide, theme: BrandTheme, page: Any,
               fitter: _Fitter) -> None:
    x, y, w, h = _bbox(c, slide.grid)
    color = theme.white if c.variant == "white" else theme.black
    _text(page, x, y, w, h, "Choose", name=c.object_id, fitter=fitter,
          font=theme.serif_family, size=_LOGO_PT[c.size], bold=True, italic=True,
          color=color)


_EMITTERS: dict[str, Callable[..., None]] = {
    "header": _emit_header,
    "subtitle": _emit_subtitle,
    "eyebrow": _emit_eyebrow,
    "body": _emit_body,
    "kpi": _emit_kpi,
    "image": _emit_image,
    "logo": _emit_logo,
    "panel": _emit_panel,
    "table": _emit_table,
}


# ---- top-level -------------------------------------------------------
def render(
    deck: Deck, *, theme: BrandTheme | None = None, fonts: Fonts | None = None
) -> Presentation:
    """Render a validated Deck into an in-memory python-pptx Presentation.

    Re-validates the deck first (overlap / out-of-bounds / duplicate ids),
    then draws each slide on a blank 16:9 layout. Returns the Presentation;
    the caller saves it (`slidebox.save`) or uploads it
    (`slidebox.to_google_slides`).

    `fonts` maps a theme family name (e.g. "Maison Neue") to a font-file path
    or a dict of `regular`/`bold`/`italic`/`bold_italic` paths. When given,
    text is sized from the real font metrics so it provably fits its box;
    otherwise a conservative width estimate is used. The files are read for
    measurement only — they need not be installed.
    """
    deck = Deck.model_validate(deck.model_dump())
    theme = theme or BrandTheme()
    fitter = _Fitter(_FontBook(fonts) if fonts else None)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank = prs.slide_layouts[6]

    for slide in deck.slides:
        page = prs.slides.add_slide(blank)
        page.background.fill.solid()
        page.background.fill.fore_color.rgb = _rgb(theme.background_rgb(slide.background))
        # Panels are background underlays: draw them first (stable order) so
        # they sit behind the image/text placed on top.
        ordered = sorted(slide.cards, key=lambda c: 0 if c.type == "panel" else 1)
        for card in ordered:
            _EMITTERS[card.type](card, slide, theme, page, fitter)

    return prs


__all__ = ["Fonts", "render"]
